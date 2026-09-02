"""
Generador de reportes: mapas + gráficos en JPG + presentación PowerPoint.

Genera:
1. Mapas coroplético por cuenca (con leyenda, escala, norte, cuadrícula)
2. Gráficos de cantidad de obras por tipo de uso (JPG)
3. Gráfico de derechos otorgados separado por cuenca (1 imagen por cuenca)
4. Presentación PowerPoint con todo integrado (tema diurno)

Uso:
    python generar_reportes.py
"""

import json
from pathlib import Path
import warnings
warnings.filterwarnings("ignore")

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.colors import Normalize, LinearSegmentedColormap
from matplotlib.gridspec import GridSpec
import geopandas as gpd
from shapely.geometry import shape
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Advertencia: python-pptx no instalado. La PPT no se generará.")

# Configuración
plt.style.use('seaborn-v0_8-whitegrid')
plt.rcParams['figure.facecolor'] = 'white'
plt.rcParams['axes.facecolor'] = '#f8f9fa'
plt.rcParams['font.family'] = 'sans-serif'
plt.rcParams['font.size'] = 10

ROOT = Path(__file__).parent
DATA_DIR = ROOT / "data"
OUTPUT_DIR = ROOT / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

TIPOS_ORDER = [
    "Represa Grande", "Represa Mediana", "Represa Chica",
    "Tajamar Grande", "Tajamar Mediano", "Tajamar Chico",
    "Tanque Excavado", "Reservorio/Otros",
]

N1_COLORS = {
    "Río Uruguay": "#3987e5",
    "Río de la Plata": "#d95926",
    "Océano Atlántico": "#199e70",
    "Laguna Merín": "#c98500",
    "Río Negro": "#d55181",
    "Santa Lucía": "#2fa72f",
}

TIPO_COLORS = {
    "Represa Grande": "#3987e5",
    "Represa Mediana": "#d95926",
    "Represa Chica": "#199e70",
    "Tajamar Grande": "#c98500",
    "Tajamar Mediano": "#d55181",
    "Tajamar Chico": "#008300",
    "Tanque Excavado": "#9085e9",
    "Reservorio/Otros": "#e66767",
}

# Rampa secuencial para mapa (azules)
SEQ_COLORS = ["#cde2fb", "#9ec5f4", "#6da7ec", "#3987e5", "#256abf", "#184f95", "#0d366b"]

def load_data():
    """Carga datos de solicitudes y cuencas."""
    df = pd.read_csv(DATA_DIR / "solicitudes_limpio.csv")
    df["codcuenca"] = df["codcuenca"].astype(int)
    
    with open(DATA_DIR / "cuencas_n2.geojson", encoding="utf-8") as f:
        geojson = json.load(f)
    
    props = pd.DataFrame([feat["properties"] for feat in geojson["features"]])
    props["codcuenca"] = props["codcuenca"].astype(int)
    props = props.rename(columns={"nombre_cue": "nombre_cuenca", "area": "area_km2"})
    
    # Crear GeoDataFrame desde GeoJSON
    gdf = gpd.GeoDataFrame.from_features(geojson["features"])
    gdf["codcuenca"] = gdf["codcuenca"].astype(int)
    gdf = gdf.rename(columns={"nombre_cue": "nombre_cuenca"})
    
    return df, geojson, props, gdf

def add_map_elements(ax, gdf):
    """Añade leyenda, escala, cuadrícula y norte a un mapa."""
    # Cuadrícula de coordenadas
    ax.gridlines(draw_labels=True, linestyle='--', linewidth=0.5, alpha=0.5, color='gray')
    
    # Escala (simple, usando coordenadas)
    ax.set_aspect('equal')
    
    # Añadir norte (flecha)
    from matplotlib.patches import FancyArrowPatch
    xlim = ax.get_xlim()
    ylim = ax.get_ylim()
    arrow_x = xlim[0] + (xlim[1] - xlim[0]) * 0.08
    arrow_y = ylim[1] - (ylim[1] - ylim[0]) * 0.08
    arrow = FancyArrowPatch(
        (arrow_x, arrow_y), (arrow_x, arrow_y + (ylim[1] - ylim[0]) * 0.05),
        arrowstyle='->', mutation_scale=20, color='black', linewidth=2
    )
    ax.add_patch(arrow)
    ax.text(arrow_x, arrow_y + (ylim[1] - ylim[0]) * 0.06, 'N', 
            ha='center', va='bottom', fontsize=12, fontweight='bold')

def generar_mapa_coplenetico(df, gdf, props):
    """Genera mapa coroplético de volumen por cuenca nivel 2."""
    print("Generando mapa coroplético...")
    
    # Agregar datos por cuenca
    por_cuenca = (
        df.groupby("codcuenca")
        .agg(volumen=("volumen", "sum"), n_obras=("volumen", "count"))
        .reset_index()
    )
    
    # Merge con geometría
    gdf_merged = gdf.merge(por_cuenca, on="codcuenca", how="left")
    gdf_merged["volumen"] = gdf_merged["volumen"].fillna(0)
    
    # Crear figura
    fig, ax = plt.subplots(figsize=(16, 10), facecolor='white')
    ax.set_facecolor('#e8f4f8')
    
    # Colormap
    vmin = gdf_merged[gdf_merged["volumen"] > 0]["volumen"].min()
    vmax = gdf_merged["volumen"].quantile(0.95)
    norm = Normalize(vmin=vmin, vmax=vmax)
    cmap = LinearSegmentedColormap.from_list("seq", SEQ_COLORS)
    
    # Plot
    gdf_merged.plot(
        column="volumen",
        ax=ax,
        legend=False,
        cmap=cmap,
        norm=norm,
        edgecolor='black',
        linewidth=0.5,
        alpha=0.8
    )
    
    # Sin límites de ejes visibles
    ax.set_xlim(gdf_merged.total_bounds[0] - 0.5, gdf_merged.total_bounds[2] + 0.5)
    ax.set_ylim(gdf_merged.total_bounds[1] - 0.5, gdf_merged.total_bounds[3] + 0.5)
    ax.axis('off')
    
    # Título y leyenda
    ax.set_title("Volumen otorgado (m³) por cuenca nivel 2\nUruguay", 
                fontsize=16, fontweight='bold', pad=20)
    
    # Barra de escala de colores
    sm = plt.cm.ScalarMappable(cmap=cmap, norm=norm)
    sm.set_array([])
    cbar = fig.colorbar(sm, ax=ax, pad=0.02, aspect=40, shrink=0.8)
    cbar.set_label("Volumen (m³)", fontsize=10, fontweight='bold')
    
    # Cuadrícula
    ax.grid(True, alpha=0.2, linestyle='--', linewidth=0.5)
    
    # Añadir norte
    xlim = ax.get_xlim()
    ylim = ax.get_ylim()
    arrow_x = xlim[0] + (xlim[1] - xlim[0]) * 0.05
    arrow_y = ylim[1] - (ylim[1] - ylim[0]) * 0.05
    from matplotlib.patches import FancyArrowPatch
    arrow = FancyArrowPatch(
        (arrow_x, arrow_y), (arrow_x, arrow_y + (ylim[1] - ylim[0]) * 0.04),
        arrowstyle='->', mutation_scale=15, color='black', linewidth=2
    )
    ax.add_patch(arrow)
    ax.text(arrow_x + 0.1, arrow_y + (ylim[1] - ylim[0]) * 0.04, 'N', 
            ha='left', va='center', fontsize=11, fontweight='bold')
    
    # Guardar
    output_path = OUTPUT_DIR / "01_mapa_coplenetico_volumen.jpg"
    fig.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    print(f"  ✓ {output_path}")
    plt.close(fig)
    
    return output_path

def generar_graficos_cantidad_obras(df):
    """Genera gráficos de cantidad de obras por tipo y uso."""
    print("Generando gráficos de cantidad de obras...")
    
    # Pivot table: tipo de obra vs uso
    pivot = df.pivot_table(
        index="tipo_obra_agr",
        columns="uso",
        values="volumen",
        aggfunc="count",
        fill_value=0
    )
    pivot = pivot.reindex(TIPOS_ORDER, fill_value=0)
    
    # Gráfico 1: Total de obras por tipo
    fig, axes = plt.subplots(1, 2, figsize=(16, 6), facecolor='white')
    
    # Barplot horizontal
    tipos = pivot.index.tolist()
    totales = pivot.sum(axis=1).values
    colors_tipos = [TIPO_COLORS.get(t, '#999999') for t in tipos]
    
    axes[0].barh(tipos, totales, color=colors_tipos, edgecolor='black', linewidth=1)
    axes[0].set_xlabel("Cantidad de obras", fontsize=12, fontweight='bold')
    axes[0].set_title("Cantidad total de obras por tipo", fontsize=13, fontweight='bold')
    axes[0].grid(axis='x', alpha=0.3)
    for i, v in enumerate(totales):
        axes[0].text(v + 10, i, str(int(v)), va='center', fontweight='bold')
    
    # Gráfico 2: Desglose por uso (stacked)
    pivot.plot(kind='barh', stacked=True, ax=axes[1], 
               color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd', '#8c564b', '#e377c2', '#7f7f7f'][:len(pivot.columns)])
    axes[1].set_xlabel("Cantidad de obras", fontsize=12, fontweight='bold')
    axes[1].set_ylabel("")
    axes[1].set_title("Desglose de obras por tipo y uso", fontsize=13, fontweight='bold')
    axes[1].legend(title="Uso", bbox_to_anchor=(1.05, 1), loc='upper left', fontsize=9)
    axes[1].grid(axis='x', alpha=0.3)
    
    plt.tight_layout()
    output_path = OUTPUT_DIR / "02_cantidad_obras_por_tipo_uso.jpg"
    fig.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    print(f"  ✓ {output_path}")
    plt.close(fig)
    
    return output_path, pivot

def generar_graficos_por_cuenca(df):
    """Genera un gráfico separado por cada cuenca nivel 2."""
    print("Generando gráficos por cuenca nivel 2...")
    
    paths = []
    for cod_cuenca in sorted(df["codcuenca"].unique()):
        df_cuenca = df[df["codcuenca"] == cod_cuenca]
        cuenca_n1 = df_cuenca["cuenca_n1"].mode()[0] if len(df_cuenca["cuenca_n1"].mode()) > 0 else "N/A"
        n_obras = len(df_cuenca)
        
        # Pivot por tipo de obra
        pivot = df_cuenca.pivot_table(
            index="tipo_obra_agr",
            columns="uso",
            values="volumen",
            aggfunc="count",
            fill_value=0
        )
        pivot = pivot.reindex(TIPOS_ORDER, fill_value=0)
        
        if pivot.empty or pivot.sum().sum() == 0:
            continue
        
        # Crear figura
        fig, ax = plt.subplots(figsize=(12, 6), facecolor='white')
        
        pivot.plot(kind='barh', stacked=True, ax=ax, 
                  color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd', '#8c564b', '#e377c2', '#7f7f7f'][:len(pivot.columns)])
        
        ax.set_xlabel("Cantidad de obras", fontsize=11, fontweight='bold')
        ax.set_ylabel("")
        ax.set_title(f"Cuenca #{cod_cuenca} | {cuenca_n1}\nDesglose de obras por tipo y uso (Total: {n_obras})", 
                    fontsize=12, fontweight='bold')
        ax.legend(title="Uso", fontsize=9, title_fontsize=9, loc='best')
        ax.grid(axis='x', alpha=0.3)
        
        plt.tight_layout()
        output_path = OUTPUT_DIR / f"03_cuenca_{cod_cuenca:03d}_derechos.jpg"
        fig.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
        print(f"  ✓ {output_path.name}")
        plt.close(fig)
        
        paths.append((cod_cuenca, cuenca_n1, output_path))
    
    return paths

def generar_ppt(mapa_path, obras_path, cuenca_paths):
    """Genera presentación PowerPoint con los gráficos."""
    if not HAS_PPTX:
        print("Advertencia: python-pptx no instalado, omitiendo PPT.")
        return None
    
    print("Generando presentación PowerPoint...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Diapositiva 1: Portada
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(57, 135, 229)  # Color primario
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Balance Hídrico por Cuenca Nivel 2"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Derechos de uso otorgados - Volúmenes y Obras\n{datetime.now().strftime('%d de %B de %Y')}"
    subtitle_frame.paragraphs[0].font.size = Pt(18)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Diapositiva 2: Mapa coroplético
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title = slide.shapes.title
    title.text = "Volumen otorgado por cuenca nivel 2"
    slide.shapes.add_picture(str(mapa_path), Inches(0.5), Inches(1), width=Inches(9))
    
    # Diapositiva 3: Gráficos de cantidad de obras
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Cantidad de obras por tipo y uso"
    slide.shapes.add_picture(str(obras_path), Inches(0.3), Inches(1), width=Inches(9.4))
    
    # Diapositivas 4+: Cuencas individuales
    for cod, n1, path in cuenca_paths[:15]:  # Limitar a 15 cuencas por deslizamiento
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"Cuenca #{cod} | {n1}"
        slide.shapes.add_picture(str(path), Inches(0.3), Inches(1), width=Inches(9.4))
    
    # Diapositiva final: Resumen
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 242, 246)
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Balance Hídrico por Cuenca Nivel 2"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(57, 135, 229)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "Análisis de derechos de uso otorgados"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)
    
    output_path = OUTPUT_DIR / "REPORTES_BALANCE_HIDRICO.pptx"
    prs.save(output_path)
    print(f"  ✓ {output_path}")
    
    return output_path

def main():
    print("=" * 70)
    print("GENERADOR DE REPORTES — Balance Hídrico por Cuenca")
    print("=" * 70)
    
    # Cargar datos
    df, geojson, props, gdf = load_data()
    print(f"✓ Datos cargados: {len(df)} solicitudes, {len(gdf)} cuencas")
    
    # Generar mapas y gráficos
    mapa_path = generar_mapa_coplenetico(df, gdf, props)
    obras_path, pivot = generar_graficos_cantidad_obras(df)
    cuenca_paths = generar_graficos_por_cuenca(df)
    
    # Generar PPT
    ppt_path = generar_ppt(mapa_path, obras_path, cuenca_paths)
    
    print("\n" + "=" * 70)
    print("✓ REPORTES GENERADOS EXITOSAMENTE")
    print("=" * 70)
    print(f"\nArchivos en: {OUTPUT_DIR}/")
    print(f"  - Mapas: *.jpg")
    print(f"  - Presentación: {ppt_path.name if ppt_path else 'N/A'}")
    print()

if __name__ == "__main__":
    main()
